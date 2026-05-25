"""Smoke testy PDF/PPTX extractoru — vytváří fixtures za běhu."""

from __future__ import annotations

from pathlib import Path

import pytest


def test_pdf_extract_with_text(tmp_path: Path) -> None:
    pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas

    pdf_path = tmp_path / "test.pdf"
    c = canvas.Canvas(str(pdf_path))
    c.drawString(100, 750, "Cesta a překazka.")
    c.showPage()
    c.drawString(100, 750, "Druhá strana s nadpisem.")
    c.save()

    from app.core.pdf_extract import extract_pdf_text

    result = extract_pdf_text(pdf_path, source_label="Test")
    assert "Cesta" in result.text
    assert "Druh" in result.text
    assert result.slide_count == 2


def test_pdf_extract_empty_scan_like(tmp_path: Path) -> None:
    pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas

    pdf_path = tmp_path / "empty.pdf"
    c = canvas.Canvas(str(pdf_path))
    c.showPage()
    c.save()

    from app.core.pdf_extract import extract_pdf_text

    result = extract_pdf_text(pdf_path, source_label="Empty")
    assert result.text == ""
    assert result.slide_count >= 1


def test_pptx_extract(tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation

    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Nadpis testovacího slidu"
    # Druhý slide bez obsahu (test že nevyhodí chybu)
    prs.slides.add_slide(prs.slide_layouts[6])  # blank
    prs.save(str(pptx_path))

    from app.core.pptx_extract import extract_pptx_text

    result = extract_pptx_text(pptx_path, source_label="Test PPTX")
    assert "Nadpis" in result.text
    assert result.slide_count == 2
