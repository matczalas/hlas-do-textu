"""Extrakce textu z PPTX prezentací přes python-pptx.

Čte všechny shape.text_frame textury. SmartArt / WordArt nečte (limit knihovny)."""

from __future__ import annotations

from pathlib import Path

from loguru import logger

from app.core.models import SlideText


class PptxExtractError(RuntimeError):
    pass


def extract_pptx_text(pptx_path: Path, source_label: str) -> SlideText:
    from pptx import Presentation  # lazy import

    pptx_path = Path(pptx_path)
    if not pptx_path.is_file():
        raise PptxExtractError(f"PPTX nenalezen: {pptx_path}")

    parts: list[str] = []
    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:
        raise PptxExtractError(f"Chyba při otevření PPTX {pptx_path.name}: {exc}") from exc

    slide_count = 0
    for idx, slide in enumerate(prs.slides, start=1):
        slide_count = idx
        slide_chunks: list[str] = []
        for shape in slide.shapes:
            # Textové rámce (běžné textboxy, nadpisy)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # paragraph.text zachytí i text v polích (datum, číslo slidu),
                    # který "".join(runs) vynechá.
                    line = paragraph.text.strip()
                    if line:
                        slide_chunks.append(line)
            # Tabulky — ve slidech přednášek časté; bez tohoto se tiše ztratí
            elif shape.has_table:
                for trow in shape.table.rows:
                    cells = [cell.text.strip() for cell in trow.cells]
                    row_text = " | ".join(c for c in cells if c)
                    if row_text:
                        slide_chunks.append(row_text)

        # Notes (speaker notes)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_chunks.append(f"[Poznámky:] {notes}")

        if slide_chunks:
            body = "\n".join(slide_chunks)
            parts.append(f"[Slide {idx}]\n{body}")

    text = "\n\n".join(parts)
    if not text:
        logger.warning("PPTX {} neobsahuje extrahovatelný text", pptx_path.name)
    else:
        logger.info("PPTX {}: {} slidů, {} znaků", pptx_path.name, slide_count, len(text))

    return SlideText(source_label=source_label, text=text, slide_count=slide_count)
